from __future__ import annotations

from pathlib import Path
import zipfile

try:
    import pypdf
except ModuleNotFoundError:  # pragma: no cover - ambiente sem dependência opcional
    pypdf = None

try:
    from docx import Document
except ModuleNotFoundError:  # pragma: no cover - ambiente sem dependência opcional
    Document = None

try:
    from openpyxl import load_workbook
except ModuleNotFoundError:  # pragma: no cover - ambiente sem dependência opcional
    load_workbook = None

try:
    from pptx import Presentation
except ModuleNotFoundError:  # pragma: no cover - ambiente sem dependência opcional
    Presentation = None

from clawai.vision.vision import vision


class DocumentReader:

    IMAGE_EXTENSIONS = {
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
        ".gif",
        ".webp",
    }

    TEXT_EXTENSIONS = {
        ".txt",
        ".md",
        ".py",
        ".json",
        ".yaml",
        ".yml",
        ".xml",
        ".sql",
        ".csv",
        ".log",
        ".ini",
        ".toml",
        ".tsx",
        ".ts",
        ".jsx",
        ".js",
        ".html",
        ".css",
    }

    def read(
        self,
        file: str | Path,
    ) -> str:

        file = Path(file)

        suffix = file.suffix.lower()

        if suffix in self.IMAGE_EXTENSIONS:
            return vision.analyze(
                file,
                "Descreva detalhadamente tudo o que existe na imagem."
            )

        if suffix == ".pdf":
            return self._pdf(file)

        if suffix == ".docx":
            return self._docx(file)

        if suffix == ".xlsx":
            return self._xlsx(file)

        if suffix == ".pptx":
            return self._pptx(file)

        if suffix == ".zip":
            return self._zip(file)

        if suffix in self.TEXT_EXTENSIONS:
            return file.read_text(
                encoding="utf-8",
                errors="ignore",
            )

        raise ValueError(
            f"Formato não suportado: {suffix}"
        )

    def _pdf(self, file: Path) -> str:
        if pypdf is None:
            raise RuntimeError("pypdf não instalado")

        reader = pypdf.PdfReader(file)

        return "\n\n".join(
            page.extract_text() or ""
            for page in reader.pages
        )

    def _docx(self, file: Path) -> str:
        if Document is None:
            raise RuntimeError("python-docx não instalado")

        doc = Document(file)

        return "\n".join(
            p.text
            for p in doc.paragraphs
            if p.text.strip()
        )

    def _xlsx(self, file: Path) -> str:
        if load_workbook is None:
            raise RuntimeError("openpyxl não instalado")

        wb = load_workbook(
            file,
            data_only=True,
        )

        lines = []

        for sheet in wb.worksheets:

            lines.append(
                f"=== PLANILHA: {sheet.title} ==="
            )

            for row in sheet.iter_rows(values_only=True):

                values = [
                    str(v)
                    for v in row
                    if v is not None
                ]

                if values:
                    lines.append(
                        " | ".join(values)
                    )

        return "\n".join(lines)

    def _pptx(self, file: Path) -> str:
        if Presentation is None:
            raise RuntimeError("python-pptx não instalado")

        prs = Presentation(file)

        lines = []

        for i, slide in enumerate(
            prs.slides,
            start=1,
        ):

            lines.append(
                f"=== SLIDE {i} ==="
            )

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text = shape.text.strip()

                    if text:
                        lines.append(text)

        return "\n".join(lines)

    def _zip(self, file: Path) -> str:

        temp = file.parent / (file.stem + "_unzipped")

        if temp.exists():

            import shutil

            shutil.rmtree(temp)

        temp.mkdir(
            parents=True,
            exist_ok=True,
        )

        with zipfile.ZipFile(file) as zip_file:
            zip_file.extractall(temp)

        files = sorted(
            p.relative_to(temp)
            for p in temp.rglob("*")
        )

        lines = []

        for item in files:

            path = temp / item

            if path.is_dir():
                continue

            lines.append(f"# {item}")

            try:
                lines.append(
                    self.read(path)
                )
            except Exception:
                lines.append(
                    "<arquivo não suportado>"
                )

            lines.append("")

        return "\n".join(lines)


documents = DocumentReader()
